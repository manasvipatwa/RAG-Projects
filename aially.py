import openai
import numpy as np
from faiss import IndexFlatL2
from transformers import pipeline, AutoTokenizer, AutoModel, AutoModelForCausalLM
import torch
import os
import requests
from opik import track, opik_context
import streamlit as st
import requests
import fitz  
import pandas as pd  
from io import StringIO
import docx 
from pptx import Presentation  
import json
from firecrawl import FirecrawlApp
from IPython.display import Markdown, display
import time

# Initialize environment variable and model
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
ollama_url = "http://localhost:11434/api/generate"  
tokenizer = AutoTokenizer.from_pretrained("sentence-transformers/all-MiniLM-L6-v2")
model = AutoModel.from_pretrained("sentence-transformers/all-MiniLM-L6-v2")
tokenizer_hf = AutoTokenizer.from_pretrained("meta-llama/Llama-3.2-1B-Instruct")
model_hf = AutoModelForCausalLM.from_pretrained("meta-llama/Llama-3.2-1B-Instruct", device_map="auto")

# Function to generate embeddings
def get_embeddings(texts):
    inputs = tokenizer(texts, padding=True, truncation=True, return_tensors="pt")
    with torch.no_grad():
        embeddings = model(**inputs).last_hidden_state.mean(dim=1).cpu().numpy()
    return embeddings

# Function to generate embeddings and store them in Faiss
def generate_embeddings(documents, file_names):
    embeddings = get_embeddings(documents)
    
    # Store embeddings in Faiss
    index = IndexFlatL2(embeddings.shape[1])
    index.add(embeddings)
    
    return embeddings, index, file_names
    

# Function to find the most relevant document based on Faiss similarity search
def find_relevant_document(question, faiss_index, documents, file_names):
    question_embedding = get_embeddings([question])
    D, I = faiss_index.search(question_embedding, k = len(documents) )
    file_name = file_names[I[0][0]]
    return documents[I[0][0]], D[0][0], file_name  


@track(project_name="Document QA HF", tags=['hugging-face', 'python-library', 'querying'])

def get_huggingface_answer(question, document):
    # Construct the prompt with the document and question
    prompt = f"Document: {document}\n\nQuestion: {question}\nAnswer:"
    
    # Track the start time to measure the duration
    start_time = time.time()
    
    # Tokenize the prompt and generate an answer
    input_ids = tokenizer_hf.encode(prompt, return_tensors="pt").to(model_hf.device)
    
    # Generate the answer
    output_ids = model_hf.generate(
        input_ids=input_ids,
        max_new_tokens=600,
        do_sample=True,
        temperature=0.1,
    )

    # Calculate the time taken for the response
    eval_duration = time.time() - start_time
    
    # Decode and return the answer
    answer = tokenizer_hf.decode(output_ids[0], skip_special_tokens=True)
    answer = answer.split("Answer:")[-1].strip()  # Clean the response
    
    # Track Opik context
    opik_context.update_current_span(
        tags=['hugging-face', 'query-answering'],  
        metadata={
            'model': "meta-llama/Llama-3.2-1B-Instruct",  
            'eval_duration': eval_duration,
            'load_duration': None,  
            'prompt_eval_duration': eval_duration,  
            'prompt_eval_count': len(input_ids[0]),
            'done': True,
            'done_reason': "Completed successfully",
        },
        usage={
            'completion_tokens': len(output_ids[0]),  # Number of tokens in the completion
            'prompt_tokens': len(input_ids[0]),  # Number of tokens in the input
            'total_tokens': len(output_ids[0]) + len(input_ids[0])  # Total tokens used
        }
    )
    
    return answer.strip()

# Streamlit frontend
st.title('AI Ally')

# Create two columns for cards
col1, col2 = st.columns(2, gap="small", vertical_alignment="center")



# Card 1: Document Upload and QA
with col1:
    if st.button("Document Insights", use_container_width=True):
        # Hide other sections and show Document Upload section
        st.session_state["section"] = "document_upload"

# Card 2: Website Scraping and QA
with col2:
    if st.button("Website Insights", use_container_width=True):
        # Hide other sections and show Website Scraping section
        st.session_state["section"] = "website_scraping"

# Handle what section to display based on button clicked
if "section" not in st.session_state:
    st.session_state["section"] = "none"

# Section for Document Upload and QA
if st.session_state["section"] == "document_upload":
    uploaded_files = st.file_uploader("Choose files", type=["pdf", "txt", "docx", "xlsx", "pptx", "csv"], accept_multiple_files=True)


    # Helper function to read text from DOCX files
    def extract_text_from_docx(docx_file):
        doc = docx.Document(docx_file)
        doc_text = ""
        for para in doc.paragraphs:
            doc_text += para.text + "\n"
        return doc_text

    # Helper function to read text from Excel files
    def extract_text_from_excel(excel_file):
        df = pd.read_excel(excel_file)
        # Concatenate all columns into a single string
        excel_text = ""
        for col in df.columns:
            excel_text += "\n".join(df[col].astype(str)) + "\n"
        return excel_text

    # Helper function to read text from PPT files
    def extract_text_from_pptx(pptx_file):
        prs = Presentation(pptx_file)
        pptx_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    pptx_text += shape.text + "\n"
        return pptx_text

    # Process each uploaded file
    if uploaded_files:
        all_documents = []
        file_names = []
        for uploaded_file in uploaded_files:
            # Process PDF files
            if uploaded_file.type == "application/pdf":
                pdf_file = fitz.open(stream=uploaded_file.read(), filetype="pdf")
                doc_text = ""
                for page in pdf_file:
                    doc_text += page.get_text()
                all_documents.append(doc_text)
                file_names.append(uploaded_file.name)
                # st.write(f'Processed PDF file: {uploaded_file.name}')
                # Provide a button to preview PDF file
                if st.button(f"Preview {uploaded_file.name}"):
                    st.text_area("Preview PDF file", doc_text[:500])
                    # display_pdf(uploaded_file)

            # Process DOCX files
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc_text = extract_text_from_docx(uploaded_file)
                all_documents.append(doc_text)
                file_names.append(uploaded_file.name)
                # st.write(f'Processed DOCX file: {uploaded_file.name}')
                # Provide a button to preview DOCX file
                if st.button(f"Preview {uploaded_file.name}"):
                    st.text_area("Preview DOCX file", doc_text[:500])

            # Process TXT files
            elif uploaded_file.type == "text/plain":
                doc_text = StringIO(uploaded_file.getvalue().decode("utf-8")).read()
                all_documents.append(doc_text)
                file_names.append(uploaded_file.name)
                # st.write(f'Processed TXT file: {uploaded_file.name}')
                # Provide a button to preview TXT file
                if st.button(f"Preview {uploaded_file.name}"):
                    st.text_area("Preview TXT file", doc_text[:500])
            
            # Process Excel files
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                doc_text = extract_text_from_excel(uploaded_file)
                all_documents.append(doc_text)
                # st.write(f'Processed Excel file: {uploaded_file.name}')
                # Provide a button to preview Excel file
                if st.button(f"Preview {uploaded_file.name}"):
                    st.text_area("Preview Excel file", doc_text[:500])
            
            # Process PPTX files
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                doc_text = extract_text_from_pptx(uploaded_file)
                all_documents.append(doc_text)
                file_names.append(uploaded_file.name)
                # st.write(f'Processed PPTX file: {uploaded_file.name}')
                # Provide a button to preview PPTX file
                if st.button(f"Preview {uploaded_file.name}"):
                    st.text_area("Preview PPTX file", doc_text[:500])
            
            # Process CSV files
            elif uploaded_file.type == "text/csv":
                df = pd.read_csv(uploaded_file)
                csv_text = df.to_string(index=False)
                all_documents.append(csv_text)
                file_names.append(uploaded_file.name)
                # st.write(f'Processed CSV file: {uploaded_file.name}')
                # Provide a button to preview CSV file
                if st.button(f"Preview {uploaded_file.name}"):
                    st.text_area("Preview CSV file", csv_text[:500])

        # After processing all files, generate embeddings for the combined documents
        embeddings, index, file_names = generate_embeddings(all_documents, file_names)
        # st.write("Embeddings generated and stored in Faiss.")
               
        example_questions = []
        for file in file_names:
            example_questions.extend([
                f"What is the main topic of {file}?",
                f"What are the key points of {file}?",
                f"Can you extract any conclusions from {file}?"
            ])
        
        # Display the predefined questions in a row
        st.subheader("Query:")
        
        # Ask a custom question
        user_question = st.text_input('Ask a question specifying the file name:')
        
        # Create columns dynamically for each question based on the number of predefined questions
        questions_per_row = 3  
        rows = (len(example_questions) + questions_per_row - 1) // questions_per_row  # Calculate the number of rows
        
        # Initialize a list to store the answers
        answers = []
        st.write('Example Questions:')
        # Display questions in rows 
        for i in range(rows):
            cols = st.columns(questions_per_row)  # Create 5 columns per row
            start_idx = i * questions_per_row
            end_idx = min((i + 1) * questions_per_row, len(example_questions))  # Ensure the last row doesn't exceed the list length
            
            # Add a button for each predefined question in the current row
            for j, predefined_question in enumerate(example_questions[start_idx:end_idx]):
                with cols[j]:
                    if st.button(predefined_question):
                        # Determine which file the question is related to based on its index
                        relevant_doc, distance, file_name = find_relevant_document(predefined_question, index, all_documents, file_names)
                        answer = get_huggingface_answer(predefined_question, relevant_doc)
                        answers.append({
                            'question': predefined_question,
                            'file': file_name,
                            'distance': distance,
                            'answer': answer
                        })
        
        # Process the custom question if provided
        if user_question:
            relevant_doc, distance, file_name = find_relevant_document(user_question, index, all_documents, file_names)
            answer = get_huggingface_answer(user_question, relevant_doc)
            answers.append({
                'question': user_question,
                'file': file_name,
                'distance': distance,
                'answer': answer
            })
        
        # Display all responses 
        if answers:
            st.subheader("Response:")
            for ans in answers:
                st.write(f"**Question:** {ans['question']}")
                st.write(f"**Found in file:** {ans['file']}")
                st.write(f"**Relevant document distance:** {ans['distance']}")
                st.write(f"**Answer:** {ans['answer']}")

# Section for Website Scraping and QA
if st.session_state["section"] == "website_scraping":
    url_input = st.text_input("Enter the URL of the website:")    
    if url_input:
        documents = []
        file_names = [url_input] 
        firecrawl_api_key = os.getenv('FIRECRAWL_API')

        # Initialize FireCrawlApp instance
        app = FirecrawlApp(api_key=firecrawl_api_key)

        try:
            # Scrape data using FireCrawl
            scrape_result = app.scrape_url(url_input, params={'formats': ['markdown', 'html']})

            # Extract content from the markdown field
            if 'markdown' in scrape_result:
                scraped_text = scrape_result['markdown']

                if scraped_text:
                    st.write(f"Preview website: {url_input}")
                    documents.append(scraped_text)
                else:
                    st.write(f"No valid content extracted for {url_input}. Please check the URL.")
            else:
                st.write(f"Error: 'markdown' key not found in the FireCrawl response. Full response shown above.")
        except Exception as e:
            st.error(f"Error while scraping: {e}")


        # Generate embeddings for the scraped content
        if documents:
            embeddings, index, file_names = generate_embeddings(documents, file_names)
            # st.write("Embeddings for scraped data generated and stored.")

            example_questions = [
                f"What is the main topic of this website?",
                f"Can you summarize this website?",
                f"What are the key points of this website?",
                f"What is the purpose of this website?",
                f"Can you extract any conclusions from this website?"
            ]

            st.subheader('Query:')
            # Ask the user for a question about the scraped content
            user_question = st.text_input('Ask a question:')
            questions_per_row = 3 
            rows = (len(example_questions) + questions_per_row - 1) // questions_per_row  # Calculate the number of rows
            
            # Initialize a list to store the answers
            answers = []
            st.write('Example Questions:') 
            # Display questions in rows
            for i in range(rows):
                cols = st.columns(questions_per_row)  # Create n columns per row
                start_idx = i * questions_per_row
                end_idx = min((i + 1) * questions_per_row, len(example_questions))  # Ensure the last row doesn't exceed the list length
                
                # Add a button for each predefined question in the current row
                for j, predefined_question in enumerate(example_questions[start_idx:end_idx]):
                    with cols[j]:
                        if st.button(predefined_question):
                            relevant_doc, distance, file_name = find_relevant_document(predefined_question, index, documents, file_names)
                            answer = get_huggingface_answer(predefined_question, relevant_doc)
                            answers.append({
                                'question': predefined_question,
                                'file': file_name,
                                'distance': distance,
                                'answer': answer
                            })
            # Process the custom question if provided
            if user_question:
                relevant_doc, distance, file_name = find_relevant_document(user_question, index, all_documents, file_names)
                answer = get_huggingface_answer(user_question, relevant_doc)
                answers.append({
                    'question': user_question,
                    'file': file_name,
                    'distance': distance,
                    'answer': answer
                })
            
            # Display responses
            if answers:
                st.subheader("Response:")
                for ans in answers:
                    st.write(f"**Question:** {ans['question']}")
                    st.write(f"**Answer:** {ans['answer']}")
        else:
                st.write("No valid scraped content to generate embeddings.")

