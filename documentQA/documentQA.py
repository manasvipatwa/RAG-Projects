import openai
import numpy as np
from faiss import IndexFlatL2
from transformers import AutoTokenizer, AutoModel
import torch
import os
import requests
from opik import track, opik_context
import streamlit as st
import requests
import fitz  
import pandas as pd  
from io import StringIO
import docx 
from pptx import Presentation  
import json

# Initialize environment variable and model
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
ollama_url = "http://localhost:11434/api/generate"  
tokenizer = AutoTokenizer.from_pretrained("sentence-transformers/all-MiniLM-L6-v2")
model = AutoModel.from_pretrained("sentence-transformers/all-MiniLM-L6-v2")

# Function to generate embeddings
def get_embeddings(texts):
    inputs = tokenizer(texts, padding=True, truncation=True, return_tensors="pt")
    with torch.no_grad():
        embeddings = model(**inputs).last_hidden_state.mean(dim=1).cpu().numpy()
    return embeddings

# Function to generate embeddings and store them in Faiss
def generate_embeddings(documents, file_names):
    embeddings = get_embeddings(documents)
    
    # Store embeddings in Faiss
    index = IndexFlatL2(embeddings.shape[1])
    index.add(embeddings)
    
    return embeddings, index, file_names

# Function to find the most relevant document based on Faiss similarity search
def find_relevant_document(question, faiss_index, documents, file_names):
    question_embedding = get_embeddings([question])
    D, I = faiss_index.search(question_embedding, k=1) 
    file_name = file_names[I[0][0]]
    return documents[I[0][0]], D[0][0], file_name  

# Track the Ollama API call 
@track(project_name="Document QA", tags=['ollama', 'python-library', 'querying'])
def get_ollama_answer(question, document):
    prompt = f"Document: {document}\n\nQuestion: {question}\nAnswer:"
    
    response = requests.post(ollama_url, json={
        "model": "llama3.2",
        "prompt": prompt,
        "stream": False  
    })

    response_data = response.json()

    # Track the response data with Opik
    if 'response' in response_data:
        answer = response_data['response'].strip() 
        
        # Update Opik context with relevant metadata 
        opik_context.update_current_span(
            tags=['ollama', 'query-answering'],  
            metadata={
                'model': response_data['model'],
                'eval_duration': response_data['eval_duration'],
                'load_duration': response_data['load_duration'],
                'prompt_eval_duration': response_data['prompt_eval_duration'],
                'prompt_eval_count': response_data['prompt_eval_count'],
                'done': response_data['done'],
                'done_reason': response_data['done_reason'],
            },
            usage={
                'completion_tokens': response_data['eval_count'],
                'prompt_tokens': response_data['prompt_eval_count'],
                'total_tokens': response_data['eval_count'] + response_data['prompt_eval_count']
            }
        )

        return answer
    else:
        raise ValueError("Error fetching answer from Ollama API")

# Streamlit frontend
st.title('Document Upload and Question Answering App')

# Upload multiple files of any type (PDF, TXT, DOCX, Excel, PPT, CSV)
uploaded_files = st.file_uploader("Choose files", type=["pdf", "txt", "docx", "xlsx", "pptx", "csv"], accept_multiple_files=True)

# Helper function to read text from DOCX files
def extract_text_from_docx(docx_file):
    doc = docx.Document(docx_file)
    doc_text = ""
    for para in doc.paragraphs:
        doc_text += para.text + "\n"
    return doc_text

# Helper function to read text from Excel files
def extract_text_from_excel(excel_file):
    df = pd.read_excel(excel_file)
    # Concatenate all columns into a single string
    excel_text = ""
    for col in df.columns:
        excel_text += "\n".join(df[col].astype(str)) + "\n"
    return excel_text

# Helper function to read text from PPT files
def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    pptx_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pptx_text += shape.text + "\n"
    return pptx_text

# Process each uploaded file
if uploaded_files:
    all_documents = []
    file_names = []
    for uploaded_file in uploaded_files:
        # Process PDF files
        if uploaded_file.type == "application/pdf":
            pdf_file = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            doc_text = ""
            for page in pdf_file:
                doc_text += page.get_text()
            all_documents.append(doc_text)
            file_names.append(uploaded_file.name)
            st.write(f'Processed PDF file: {uploaded_file.name}')

        # Process DOCX files
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc_text = extract_text_from_docx(uploaded_file)
            all_documents.append(doc_text)
            file_names.append(uploaded_file.name)
            st.write(f'Processed DOCX file: {uploaded_file.name}')

        # Process TXT files
        elif uploaded_file.type == "text/plain":
            doc_text = StringIO(uploaded_file.getvalue().decode("utf-8")).read()
            all_documents.append(doc_text)
            file_names.append(uploaded_file.name)
            st.write(f'Processed TXT file: {uploaded_file.name}')
        
        # Process Excel files
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            doc_text = extract_text_from_excel(uploaded_file)
            all_documents.append(doc_text)
            st.write(f'Processed Excel file: {uploaded_file.name}')
        
        # Process PPTX files
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            doc_text = extract_text_from_pptx(uploaded_file)
            all_documents.append(doc_text)
            file_names.append(uploaded_file.name)
            st.write(f'Processed PPTX file: {uploaded_file.name}')
        
        # Process CSV files
        elif uploaded_file.type == "text/csv":
            df = pd.read_csv(uploaded_file)
            csv_text = df.to_string(index=False)
            all_documents.append(csv_text)
            file_names.append(uploaded_file.name)
            st.write(f'Processed CSV file: {uploaded_file.name}')

    # After processing all files, generate embeddings for the combined documents
    embeddings, index, file_names = generate_embeddings(all_documents, file_names)
    st.write("Embeddings generated and stored in Faiss.")

    # User question input
    question = st.text_input('Ask a question:')
    if question:
        # Find relevant document using Faiss 
        relevant_doc, distance, file_name = find_relevant_document(question, index, all_documents, file_names)
        answer = get_ollama_answer(question, relevant_doc)
        st.write(f'Found in file: {file_name}')
        st.write(f'Relevant document distance: {distance}')
        st.write(f'Answer: {answer}')
