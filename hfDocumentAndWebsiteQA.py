import openai
import numpy as np
from faiss import IndexFlatL2
from transformers import pipeline, AutoTokenizer, AutoModel
import torch
import os
import requests
from opik import track, opik_context
import streamlit as st
import requests
import fitz  
import pandas as pd  
from io import StringIO
import docx 
from pptx import Presentation  
import json
from firecrawl import FirecrawlApp
from IPython.display import Markdown, display

# Initialize environment variable and model
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
ollama_url = "http://localhost:11434/api/generate"  
tokenizer = AutoTokenizer.from_pretrained("sentence-transformers/all-MiniLM-L6-v2")
model = AutoModel.from_pretrained("sentence-transformers/all-MiniLM-L6-v2")
model_id = "meta-llama/Llama-3.2-1B-Instruct"
pipe = pipeline(
    "text-generation",
    model=model_id,
    torch_dtype=torch.bfloat16,
    device_map="auto",
)
# Function to generate embeddings
def get_embeddings(texts):
    inputs = tokenizer(texts, padding=True, truncation=True, return_tensors="pt")
    with torch.no_grad():
        embeddings = model(**inputs).last_hidden_state.mean(dim=1).cpu().numpy()
    return embeddings

# Function to generate embeddings and store them in Faiss
def generate_embeddings(documents, file_names):
    embeddings = get_embeddings(documents)
    
    # Store embeddings in Faiss
    index = IndexFlatL2(embeddings.shape[1])
    index.add(embeddings)
    
    return embeddings, index, file_names
    

# Function to find the most relevant document based on Faiss similarity search
def find_relevant_document(question, faiss_index, documents, file_names):
    question_embedding = get_embeddings([question])
    D, I = faiss_index.search(question_embedding, k = len(documents) )
    file_name = file_names[I[0][0]]
    return documents[I[0][0]], D[0][0], file_name  


@track(project_name="Document QA HF", tags=['hugging-face', 'python-library', 'querying'])
# Function to generate answer using Hugging Face model
def get_huggingface_answer(question, document):
    # Construct the prompt to pass to the model
    prompt = f"Document: {document}\n\nQuestion: {question}\nAnswer:"
    
    # Use the Hugging Face pipeline to generate a response
    response = pipe(prompt, max_new_tokens=600, num_return_sequences=1)

    # Extract the generated answer from the response
    answer = response[0]['generated_text'].strip()
    
    return answer

# Track the Ollama API call 
@track(project_name="Document QA", tags=['ollama', 'python-library', 'querying'])
def get_ollama_answer(question, document):
    prompt = f"Document: {document}\n\nQuestion: {question}\nAnswer:"
    
    response = requests.post(ollama_url, json={
        "model": "llama3.2",
        "prompt": prompt,
        "stream": False  
    })

    response_data = response.json()

    # Track the response data with Opik
    if 'response' in response_data:
        answer = response_data['response'].strip() 
        
        # Update Opik context with relevant metadata 
        opik_context.update_current_span(
            tags=['ollama', 'query-answering'],  
            metadata={
                'model': response_data['model'],
                'eval_duration': response_data['eval_duration'],
                'load_duration': response_data['load_duration'],
                'prompt_eval_duration': response_data['prompt_eval_duration'],
                'prompt_eval_count': response_data['prompt_eval_count'],
                'done': response_data['done'],
                'done_reason': response_data['done_reason'],
            },
            usage={
                'completion_tokens': response_data['eval_count'],
                'prompt_tokens': response_data['prompt_eval_count'],
                'total_tokens': response_data['eval_count'] + response_data['prompt_eval_count']
            }
        )

        return answer
    else:
        raise ValueError("Error fetching answer from Ollama API")



# Streamlit frontend
st.title('AI Ally')

# Create two columns for cards
col1, col2 = st.columns(2, gap="small", vertical_alignment="center")



# Card 1: Document Upload and QA
with col1:
    if st.button("Document Insights", use_container_width=True):
        # Hide other sections and show Document Upload section
        st.session_state["section"] = "document_upload"

# Card 2: Website Scraping and QA
with col2:
    if st.button("Website Insights", use_container_width=True):
        # Hide other sections and show Website Scraping section
        st.session_state["section"] = "website_scraping"

# Handle what section to display based on button clicked
if "section" not in st.session_state:
    st.session_state["section"] = "none"

# # Option to choose between uploading files or scraping a website
# option = st.radio("Select an option", ("Upload Files", "Scrape Website"))

# Option 1: File Upload
# if option == "Upload Files":
#     uploaded_files = st.file_uploader("Choose files", type=["pdf", "txt", "docx", "xlsx", "pptx", "csv"], accept_multiple_files=True)

# Section for Document Upload and QA
if st.session_state["section"] == "document_upload":
    uploaded_files = st.file_uploader("Choose files", type=["pdf", "txt", "docx", "xlsx", "pptx", "csv"], accept_multiple_files=True)


    # Helper function to read text from DOCX files
    def extract_text_from_docx(docx_file):
        doc = docx.Document(docx_file)
        doc_text = ""
        for para in doc.paragraphs:
            doc_text += para.text + "\n"
        return doc_text

    # Helper function to read text from Excel files
    def extract_text_from_excel(excel_file):
        df = pd.read_excel(excel_file)
        # Concatenate all columns into a single string
        excel_text = ""
        for col in df.columns:
            excel_text += "\n".join(df[col].astype(str)) + "\n"
        return excel_text

    # Helper function to read text from PPT files
    def extract_text_from_pptx(pptx_file):
        prs = Presentation(pptx_file)
        pptx_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    pptx_text += shape.text + "\n"
        return pptx_text

    # Process each uploaded file
    if uploaded_files:
        all_documents = []
        file_names = []
        for uploaded_file in uploaded_files:
            # Process PDF files
            if uploaded_file.type == "application/pdf":
                pdf_file = fitz.open(stream=uploaded_file.read(), filetype="pdf")
                doc_text = ""
                for page in pdf_file:
                    doc_text += page.get_text()
                all_documents.append(doc_text)
                file_names.append(uploaded_file.name)
                # st.write(f'Processed PDF file: {uploaded_file.name}')
                # Provide a button to preview PDF file
                if st.button(f"Preview {uploaded_file.name}"):
                    st.text_area("Preview PDF file", doc_text[:500])
                    # display_pdf(uploaded_file)

            # Process DOCX files
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc_text = extract_text_from_docx(uploaded_file)
                all_documents.append(doc_text)
                file_names.append(uploaded_file.name)
                # st.write(f'Processed DOCX file: {uploaded_file.name}')
                # Provide a button to preview DOCX file
                if st.button(f"Preview {uploaded_file.name}"):
                    st.text_area("Preview DOCX file", doc_text[:500])

            # Process TXT files
            elif uploaded_file.type == "text/plain":
                doc_text = StringIO(uploaded_file.getvalue().decode("utf-8")).read()
                all_documents.append(doc_text)
                file_names.append(uploaded_file.name)
                # st.write(f'Processed TXT file: {uploaded_file.name}')
                # Provide a button to preview TXT file
                if st.button(f"Preview {uploaded_file.name}"):
                    st.text_area("Preview TXT file", doc_text[:500])
            
            # Process Excel files
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                doc_text = extract_text_from_excel(uploaded_file)
                all_documents.append(doc_text)
                # st.write(f'Processed Excel file: {uploaded_file.name}')
                # Provide a button to preview Excel file
                if st.button(f"Preview {uploaded_file.name}"):
                    st.text_area("Preview Excel file", doc_text[:500])
            
            # Process PPTX files
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                doc_text = extract_text_from_pptx(uploaded_file)
                all_documents.append(doc_text)
                file_names.append(uploaded_file.name)
                # st.write(f'Processed PPTX file: {uploaded_file.name}')
                # Provide a button to preview PPTX file
                if st.button(f"Preview {uploaded_file.name}"):
                    st.text_area("Preview PPTX file", doc_text[:500])
            
            # Process CSV files
            elif uploaded_file.type == "text/csv":
                df = pd.read_csv(uploaded_file)
                csv_text = df.to_string(index=False)
                all_documents.append(csv_text)
                file_names.append(uploaded_file.name)
                # st.write(f'Processed CSV file: {uploaded_file.name}')
                # Provide a button to preview CSV file
                if st.button(f"Preview {uploaded_file.name}"):
                    st.text_area("Preview CSV file", csv_text[:500])

        # After processing all files, generate embeddings for the combined documents
        embeddings, index, file_names = generate_embeddings(all_documents, file_names)
        # st.write("Embeddings generated and stored in Faiss.")

        # User question input
        st.subheader('Query:')
        question = st.text_input('Ask a question specifying the file name:')
        if question:
            # Find relevant document using Faiss 
            relevant_doc, distance, file_name = find_relevant_document(question, index, all_documents, file_names)
            answer = get_huggingface_answer(question, relevant_doc)
            st.write(f'Found in file: {file_name}')
            st.write(f'Relevant document distance: {distance}')
            st.write(f'Answer: {answer}')

# # Option 2: Scrape Website
# if option == "Scrape Website":
#     url_input = st.text_input("Enter the URL of the website to scrape:")


# Section for Website Scraping and QA
if st.session_state["section"] == "website_scraping":
    url_input = st.text_input("Enter the URL of the website:")    
    if url_input:
        documents = []
        file_names = [url_input] 
        # Load API key from config.json
        with open('config.json') as fd:
            conf = json.loads(fd.read())
            firecrawl_api_key = conf["firecrawl_api_key"]

        # Initialize FireCrawlApp instance
        app = FirecrawlApp(api_key=firecrawl_api_key)

        try:
            # Scrape data using FireCrawl
            scrape_result = app.scrape_url(url_input, params={'formats': ['markdown', 'html']})

            # Extract content from the markdown field
            if 'markdown' in scrape_result:
                scraped_text = scrape_result['markdown']

                if scraped_text:
                    st.write(f"Scraped data from {url_input}")
                    documents.append(scraped_text)
                else:
                    st.write(f"No valid content extracted for {url_input}. Please check the URL.")
            else:
                st.write(f"Error: 'markdown' key not found in the FireCrawl response. Full response shown above.")
        except Exception as e:
            st.error(f"Error while scraping: {e}")

        # Generate embeddings for the scraped content
        if documents:
            embeddings, index, file_names = generate_embeddings(documents, file_names)
            st.write("Embeddings for scraped data generated and stored.")

            # Ask the user for a question about the scraped content
            question = st.text_input('Ask a question:')

            if question:
                # Generate relevant insights using Faiss
                relevant_doc, distance, file_name = find_relevant_document(question, index, documents, file_names)

                # Query Ollama for insights based on the relevant document
                answer = get_ollama_answer(question, relevant_doc)
                st.write(f'Answer: {answer}')
        else:
            st.write("No valid scraped content to generate embeddings.")
